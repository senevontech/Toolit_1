#!/usr/bin/env python3

import sys
import os
import shutil
import subprocess
import tempfile


# ---------------- PDF → WORD ----------------
def pdf_to_docx(input_path, output_path):
    from pdf2docx import Converter
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()


# ---------------- PDF → EXCEL ----------------
def pdf_to_xlsx(input_path, output_path):
    import pdfplumber
    import openpyxl
    import re

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    wrote_anything = False

    def write_row(ws, row_idx, values):
        nonlocal wrote_anything
        for col_idx, cell in enumerate(values, start=1):
            value = "" if cell is None else str(cell)
            if value.strip():
                wrote_anything = True
            ws.cell(row=row_idx, column=col_idx, value=value)

    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            ws = wb.create_sheet(title=f"Page {page_num}")

            tables = page.extract_tables()
            if tables:
                row_cursor = 1
                for table in tables:
                    for row in table:
                        write_row(ws, row_cursor, row)
                        row_cursor += 1
                    row_cursor += 1
            else:
                text = page.extract_text() or ""
                for row_idx, line in enumerate(text.split("\n"), start=1):
                    columns = [part for part in re.split(r"\s{2,}|\t+", line.strip()) if part]
                    write_row(ws, row_idx, columns or [line])

            for column_cells in ws.columns:
                max_length = max(len(str(cell.value or "")) for cell in column_cells)
                ws.column_dimensions[column_cells[0].column_letter].width = min(max(max_length + 2, 12), 60)

    if not wrote_anything:
        ws = wb.create_sheet(title="No extractable data")
        ws.cell(
            row=1,
            column=1,
            value="No selectable text or tables were found in this PDF. Scanned/image-only PDFs need OCR before Excel conversion.",
        )

    wb.save(output_path)


# ---------------- PDF → PPT ----------------
def pdf_to_pptx(input_path, output_path):
    from pptx import Presentation
    import tempfile

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    with tempfile.TemporaryDirectory() as tmp_dir:
        image_paths = render_pdf_pages(input_path, tmp_dir)
        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, 0, 0, slide_width, slide_height)

    prs.save(output_path)


def render_pdf_pages(input_path, tmp_dir, dpi=150):
    try:
        import fitz

        doc = fitz.open(input_path)
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)
        image_paths = []

        for idx, page in enumerate(doc):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            img_path = os.path.join(tmp_dir, f"page_{idx}.png")
            pixmap.save(img_path)
            image_paths.append(img_path)

        doc.close()

        if not image_paths:
            raise Exception("PDF did not contain any renderable pages.")

        return image_paths
    except ImportError:
        pass

    try:
        from pdf2image import convert_from_path

        images = convert_from_path(input_path, dpi=dpi)
        image_paths = []
        for idx, image in enumerate(images):
            img_path = os.path.join(tmp_dir, f"page_{idx}.png")
            image.save(img_path, "PNG")
            image_paths.append(img_path)

        if not image_paths:
            raise Exception("PDF did not contain any renderable pages.")

        return image_paths
    except Exception as exc:
        raise Exception(
            "Could not render PDF pages. Install PyMuPDF, or install Poppler and add it to PATH. "
            f"Details: {exc}"
        )


# ---------------- 🔐 PROTECT PDF ----------------
def pdf_protect(input_path, output_path, password):
    from PyPDF2 import PdfReader, PdfWriter

    reader = PdfReader(input_path)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(password)

    with open(output_path, "wb") as f:
        writer.write(f)


# ---------------- 🔓 UNLOCK PDF ----------------
def pdf_unlock(input_path, output_path, password):
    from PyPDF2 import PdfReader, PdfWriter

    reader = PdfReader(input_path)

    if reader.is_encrypted:
        result = reader.decrypt(password)
        if result == 0:
            raise Exception("Wrong password")

    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)


# ---------------- OFFICE -> PDF / HTML ----------------
def find_soffice():
    env_path = os.environ.get("LIBREOFFICE_PATH")
    candidates = []

    if env_path:
        candidates.append(env_path)

    candidates.extend([
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files\LibreOffice\program\soffice.com",
        r"C:\LibreOffice\program\soffice.exe",
        r"C:\LibreOffice\program\soffice.com",
    ])

    for candidate in candidates:
        if candidate and os.path.exists(candidate):
            return candidate

    raise Exception("LibreOffice was not found. Install LibreOffice or set LIBREOFFICE_PATH.")


def office_convert(input_path, output_path, target_format):
    soffice = find_soffice()
    source_name = os.path.splitext(os.path.basename(input_path))[0]

    with tempfile.TemporaryDirectory() as tmp_dir:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--norestore",
                "--nodefault",
                "--nofirststartwizard",
                "--nolockcheck",
                "--convert-to",
                target_format,
                "--outdir",
                tmp_dir,
                input_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode != 0:
            raise Exception((result.stderr or result.stdout or "LibreOffice conversion failed").strip())

        expected_ext = "html" if target_format == "html" else target_format
        expected_path = os.path.join(tmp_dir, f"{source_name}.{expected_ext}")

        if not os.path.exists(expected_path):
            matches = [
                os.path.join(tmp_dir, name)
                for name in os.listdir(tmp_dir)
                if name.lower().endswith(f".{expected_ext}")
            ]
            if not matches:
                raise Exception("LibreOffice did not generate an output file.")
            expected_path = matches[0]

        shutil.copyfile(expected_path, output_path)


# ---------------- MAIN ----------------
def main():
    if len(sys.argv) < 4:
        print("Usage: convert_pdf.py <mode> <input> <output> [password]", file=sys.stderr)
        sys.exit(1)

    mode = sys.argv[1].lower()
    input_path = sys.argv[2]
    output_path = sys.argv[3]

    if not os.path.isfile(input_path):
        print(f"Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    try:
        if mode == "docx":
            pdf_to_docx(input_path, output_path)

        elif mode == "xlsx":
            pdf_to_xlsx(input_path, output_path)

        elif mode == "pptx":
            pdf_to_pptx(input_path, output_path)

        elif mode == "protect":
            if len(sys.argv) != 5:
                print("Password required", file=sys.stderr)
                sys.exit(1)
            password = sys.argv[4]
            pdf_protect(input_path, output_path, password)

        elif mode == "unlock":
            if len(sys.argv) != 5:
                print("Password required", file=sys.stderr)
                sys.exit(1)
            password = sys.argv[4]
            pdf_unlock(input_path, output_path, password)

        elif mode == "excel_pdf":
            office_convert(input_path, output_path, "pdf")

        elif mode == "ppt_pdf":
            office_convert(input_path, output_path, "pdf")

        elif mode == "word_pdf":
            office_convert(input_path, output_path, "pdf")

        elif mode == "word_html":
            office_convert(input_path, output_path, "html")

        else:
            print(f"Unknown mode: {mode}", file=sys.stderr)
            sys.exit(1)

        print(f"OK:{output_path}")

    except Exception as e:
        print(f"ERROR:{str(e)}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
